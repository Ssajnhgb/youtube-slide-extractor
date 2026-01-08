import streamlit as st
import cv2
import os
import yt_dlp
from pptx import Presentation
import imagehash
from PIL import Image
import tempfile

# Set Page Config
st.set_page_config(page_title="Lecture Slide Extractor", page_icon="📊")
st.title("🎥 YouTube & Video Slide Extractor")
st.markdown("Convert video lectures into PowerPoint slides automatically.")

# Sidebar for Settings
st.sidebar.header("Extraction Settings")
threshold = st.sidebar.slider("Sensitivity (Higher = more slides)", 10, 30, 15)
frame_skip = st.sidebar.slider("Check every X frames", 30, 120, 60)

# Input Section
option = st.radio("Step 1: Choose your source", ("YouTube Link", "Upload MP4 File"))

video_file_path = None

if option == "YouTube Link":
    url = st.text_input("Paste YouTube URL here:")
    if st.button("Start Download"):
        if url:
            try:
                with st.spinner("Bypassing YouTube security..."):
                    ydl_opts = {
                        'format': 'best[ext=mp4]',
                        'outtmpl': 'downloaded_video.mp4',
                        'quiet': True,
                        'no_warnings': True,
                        'user_agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                    }
                    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                        ydl.download([url])
                    video_file_path = 'downloaded_video.mp4'
                    st.success("Download Successful!")
            except Exception as e:
                st.error(f"YouTube Blocked this server (403 Error). Please download the video to your PC and use the 'Upload MP4' option instead.")
else:
    uploaded_file = st.file_uploader("Upload a lecture video", type=["mp4", "mov", "avi"])
    if uploaded_file:
        tfile = tempfile.NamedTemporaryFile(delete=False)
        tfile.write(uploaded_file.read())
        video_file_path = tfile.name

# Processing Section
if video_file_path and st.button("Step 2: Extract Slides to PPTX"):
    try:
        cap = cv2.VideoCapture(video_file_path)
        prs = Presentation()
        last_hash = None
        count = 0
        slides_saved = 0
        
        progress_bar = st.progress(0)
        status_text = st.empty()
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))

        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break
            
            # Skip frames to speed up processing
            if count % frame_skip == 0:
                # Convert frame for hashing
                img_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                pil_img = Image.fromarray(img_rgb)
                curr_hash = imagehash.phash(pil_img)
                
                # If slide has changed significantly
                if last_hash is None or curr_hash - last_hash > threshold:
                    # Save frame to slide
                    img_path = f"frame_{slides_saved}.jpg"
                    cv2.imwrite(img_path, frame)
                    
                    slide_layout = prs.slide_layouts[6] # Blank layout
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width)
                    
                    last_hash = curr_hash
                    slides_saved += 1
                    os.remove(img_path)
                
                # Update UI
                progress = count / total_frames
                progress_bar.progress(min(progress, 1.0))
                status_text.text(f"Processing... Found {slides_saved} unique slides.")
            
            count += 1

        cap.release()
        output_pptx = "lecture_presentation.pptx"
        prs.save(output_pptx)
        
        st.success(f"Done! {slides_saved} slides extracted.")
        with open(output_pptx, "rb") as f:
            st.download_button("📥 Download PowerPoint", f, file_name="lecture_slides.pptx")
            
        # Final Cleanup
        if os.path.exists('downloaded_video.mp4'): os.remove('downloaded_video.mp4')

    except Exception as e:
        st.error(f"Critical Error: {e}")
